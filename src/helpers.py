import logging
import io
import json
from pptx import Presentation
from pypdf import PdfReader
import pypdfium2 as pdfium


# my_logger.debug('This is a debug message.')
# my_logger.info('This is an informational message.')
# my_logger.warning('This is a warning message.')
# my_logger.error('This is an error message.')
# my_logger.critical('This is a critical message!')


def json_to_dict(json_string):
    try:
        data_dict = json.loads(json_string)
        return data_dict
    except json.JSONDecodeError as e:
        print(f"Error decoding JSON: {e}")
        return None


# Load the presentation
def extract_content_from_pptx(file):
    prs = Presentation(file)
    content = {"text": [], "tables": [], "images": []}

    for slide in prs.slides:
        for shape in slide.shapes:
            # Text
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:  # type: ignore
                    content["text"].append(p.text)

            # Tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:  # type: ignore
                    table_data.append([cell.text for cell in row.cells])
                content["tables"].append(table_data)

            # Images
            if shape.shape_type == 13:  # Picture type
                image = shape.image  # type: ignore
                image_bytes = image.blob
                ext = image.ext
                content["images"].append((image_bytes, ext))

    return content


def extract_content_from_pdf(file):
    try:
        reader = PdfReader(file)
        content = {"text": [], "images": []}

        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                content["text"].append(page_text)
            for image_file_object in page.images:
                content["images"].append(image_file_object.data)

        # Fallback: if pypdf found no images, use pypdfium2 to render pages as images
        if not content["images"]:
            try:
                file.seek(0)  # reset file pointer
                pdf = pdfium.PdfDocument(file)
                for i in range(len(pdf)):
                    page = pdf[i]
                    bitmap = page.render(scale=2)  # 2x for better quality
                    pil_image = bitmap.to_pil()
                    buf = io.BytesIO()
                    pil_image.save(buf, format="PNG")
                    content["images"].append(buf.getvalue())
                pdf.close()
            except Exception as e:
                my_logger.error(f"Error extracting images with pypdfium2: {e}")

        return content
    except Exception as e:
        my_logger.error(f"Error extracting content from PDF: {e}")
        return {"text": [], "images": []}


def list_to_html_ol(cell):
    if isinstance(cell, list):
        return "<ul>" + "".join(f"<li>{item}</li>" for item in cell) + "</ul>"
    return cell


def setup_logger(name, log_file, level=logging.INFO):

    logger = logging.getLogger(name)
    logger.setLevel(level)

    file_handler = logging.FileHandler(log_file)
    file_handler.setLevel(level)

    console_handler = logging.StreamHandler()
    console_handler.setLevel(level)

    formatter = logging.Formatter(
        "%(asctime)s - %(name)s - %(levelname)s - %(message)s"
    )
    file_handler.setFormatter(formatter)
    console_handler.setFormatter(formatter)

    if not logger.handlers:
        logger.addHandler(file_handler)
        logger.addHandler(console_handler)

    return logger


my_logger = setup_logger("app_logger", "app.log")
